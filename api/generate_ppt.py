#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Amazon Sales Prediction System 0th review
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
from pathlib import Path

# Constants
SLIDE_WIDTH = Inches(13.333)  # 16:9 widescreen
SLIDE_HEIGHT = Inches(7.5)

# Color scheme
PRIMARY_COLOR = RGBColor(0x1E, 0x3A, 0x5F)  # Dark blue
SECONDARY_COLOR = RGBColor(0x34, 0x98, 0xDB)  # Light blue
ACCENT_COLOR = RGBColor(0xE7, 0x4C, 0x3C)  # Red
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)  # Light gray background
TEXT_DARK = RGBColor(0x2C, 0x3E, 0x50)  # Dark text
TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)  # White text

def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background rectangle
    left = top = 0
    width = SLIDE_WIDTH
    height = SLIDE_HEIGHT
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()
    
    # Add title text
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(11.333)
    height = Inches(1.5)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle text
    left = Inches(1)
    top = Inches(4.2)
    width = Inches(11.333)
    height = Inches(1)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER
    
    # Add presenter info
    left = Inches(1)
    top = Inches(5.5)
    width = Inches(11.333)
    height = Inches(0.5)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = "0th Review Presentation"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)  # Light gray
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content, layout=1):
    """Add a content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background
    left = top = 0
    width = SLIDE_WIDTH
    height = SLIDE_HEIGHT
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.fill.background()
    
    # Add title bar
    left = Inches(0)
    top = Inches(0)
    width = SLIDE_WIDTH
    height = Inches(1.2)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()
    
    # Add title text
    left = Inches(0.5)
    top = Inches(0.2)
    width = Inches(12.333)
    height = Inches(0.8)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    
    # Add content
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.333)
    height = Inches(5.5)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(12)
        p.space_before = Pt(6)
    
    return slide

def add_two_column_slide(prs, title, left_content, right_content, left_title="Left", right_title="Right"):
    """Add a slide with two columns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background
    left = top = 0
    width = SLIDE_WIDTH
    height = SLIDE_HEIGHT
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.fill.background()
    
    # Add title bar
    left = Inches(0)
    top = Inches(0)
    width = SLIDE_WIDTH
    height = Inches(1.2)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()
    
    # Add title text
    left = Inches(0.5)
    top = Inches(0.2)
    width = Inches(12.333)
    height = Inches(0.8)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    
    # Left column
    left_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.5))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    left_box.line.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
    
    # Left title
    left_title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.5), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Left content
    left_content_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(5.5), Inches(4.5))
    tf = left_content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_shape(1, Inches(6.9), Inches(1.5), Inches(5.9), Inches(5.5))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    right_box.line.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
    
    # Right title
    right_title_box = slide.shapes.add_textbox(Inches(7.1), Inches(1.6), Inches(5.5), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Right content
    right_content_box = slide.shapes.add_textbox(Inches(7.1), Inches(2.2), Inches(5.5), Inches(4.5))
    tf = right_content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background
    left = top = 0
    width = SLIDE_WIDTH
    height = SLIDE_HEIGHT
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.fill.background()
    
    # Add title bar
    left = Inches(0)
    top = Inches(0)
    width = SLIDE_WIDTH
    height = Inches(1.2)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()
    
    # Add title text
    left = Inches(0.5)
    top = Inches(0.2)
    width = Inches(12.333)
    height = Inches(0.8)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    
    # Add table
    table_left = Inches(0.5)
    table_top = Inches(1.5)
    table_width = Inches(12.333)
    table_height = Inches(5.5)
    
    table = slide.shapes.add_table(len(rows) + 1, len(headers), table_left, table_top, table_width, table_height).table
    
    # Set column widths
    col_width = int(table_width / len(headers))
    for i in range(len(headers)):
        table.columns[i].width = col_width
    
    # Add header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = SECONDARY_COLOR
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER
    
    # Add data rows
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_DARK
            p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_image_slide(prs, title, image_path, caption=""):
    """Add a slide with an image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background
    left = top = 0
    width = SLIDE_WIDTH
    height = SLIDE_HEIGHT
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.fill.background()
    
    # Add title bar
    left = Inches(0)
    top = Inches(0)
    width = SLIDE_WIDTH
    height = Inches(1.2)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()
    
    # Add title text
    left = Inches(0.5)
    top = Inches(0.2)
    width = Inches(12.333)
    height = Inches(0.8)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    
    # Add image if exists
    if os.path.exists(image_path):
        # Calculate image position to center it
        img_left = Inches(2)
        img_top = Inches(1.5)
        img_width = Inches(9)
        img_height = Inches(5)
        
        slide.shapes.add_picture(image_path, img_left, img_top, img_width, img_height)
        
        # Add caption if provided
        if caption:
            caption_left = Inches(1)
            caption_top = Inches(6.5)
            caption_width = Inches(11.333)
            caption_height = Inches(0.5)
            text_box = slide.shapes.add_textbox(caption_left, caption_top, caption_width, caption_height)
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = caption
            p.font.size = Pt(14)
            p.font.italic = True
            p.font.color.rgb = TEXT_DARK
            p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_presentation():
    """Create the main presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # Get screenshot paths - look for numbered files first, then any PNGs
    project_root = Path(__file__).parent.parent
    screenshots_dir = project_root / "screenshots"
    
    # Prefer specific numbered screenshots
    preferred_screenshots = [
        screenshots_dir / "1_dashboard.png",
        screenshots_dir / "2_products.png",
        screenshots_dir / "3_predictions.png"
    ]
    screenshot_files = [f for f in preferred_screenshots if f.exists()]
    
    # Fall back to any PNGs if specific files not found
    if not screenshot_files and screenshots_dir.exists():
        screenshot_files = sorted(screenshots_dir.glob("*.png"), key=lambda f: f.stat().st_mtime, reverse=True)[:3]
    
    # Slide 1: Title
    add_title_slide(prs, "Amazon Sales Prediction System", "Machine Learning for E-commerce Sales Forecasting")
    
    # Slide 2: Project Overview
    add_content_slide(prs, "Project Overview", [
        "Comprehensive ML system for predicting Amazon product sales",
        "Uses historical data with gradient boosting models (XGBoost, LightGBM)",
        "Time-series forecasting with confidence intervals",
        "Interactive web dashboard for visualization and predictions",
        "Supports what-if analysis for price/discount scenarios",
        "Built with FastAPI for low-latency predictions",
        "Deployed on Render with Vercel serverless support"
    ])
    
    # Slide 3: System Architecture
    add_two_column_slide(prs, "System Architecture", [
        "Data Layer: Synthetic data generation with realistic patterns",
        "Feature Engineering: 100+ predictive features",
        "ML Models: XGBoost, LightGBM, Random Forest ensemble",
        "Prediction Engine: Statistical fallback + ML-powered forecasts",
        "Web Interface: FastAPI with Jinja2 templates",
        "Deployment: Render (production) + Vercel (serverless)"
    ], [
        "Data Flow: Raw data → Feature engineering → Model training → Predictions",
        "Ensemble approach combines multiple models with optimized weights",
        "Time-series cross-validation prevents data leakage",
        "Lazy loading for fast startup on Render free tier",
        "Background training with progress tracking",
        "Pre-trained models for instant predictions"
    ], "Components", "Architecture")
    
    # Slide 4: Data Generation
    add_content_slide(prs, "Data Generation & Features", [
        "Generates realistic synthetic Amazon product data",
        "2,000+ products across 10 categories (Electronics, Books, etc.)",
        "365 days of daily sales history with realistic patterns:",
        "  • Power-law sales distribution (few products sell a lot)",
        "  • Price elasticity (discounts boost sales)",
        "  • Seasonality (holidays, day-of-week effects)",
        "  • Rating dynamics (drift over time with new reviews)",
        "Caches data to Parquet for fast reloading",
        "Supports real Amazon data via scraper (optional)"
    ])
    
    # Slide 5: Feature Engineering
    add_content_slide(prs, "Feature Engineering (100+ Features)", [
        "Temporal Features: Day of week, month, quarter, holiday proximity",
        "Lag Features: Sales from 1, 3, 7, 14, 30, 60, 90 days ago",
        "Rolling Statistics: 7/14/30-day moving average, std dev, min, max",
        "Price Features: Discount %, price vs category avg, price change rate",
        "Rating Features: Rating momentum, review growth rate",
        "Cross-Product Features: Category average sales, price, rating",
        "Interaction Features: Price × rating, discount × reviews, sales per dollar",
        "Advanced: Coefficient of variation, trend analysis, holiday proximity"
    ])
    
    # Slide 6: ML Models
    add_two_column_slide(prs, "Machine Learning Models", [
        "XGBoost Regressor:",
        "  • Gradient boosting with tree-based learners",
        "  • Hyperparameter tuning via Optuna",
        "  • Early stopping to prevent overfitting",
        "",
        "LightGBM Regressor:",
        "  • Leaf-wise tree growth with depth constraints",
        "  • Optimized for speed and memory efficiency"
    ], [
        "Random Forest:",
        "  • Ensemble of decision trees",
        "  • Bagging with feature randomization",
        "  • Robust baseline model",
        "",
        "Ensemble Approach:",
        "  • Combines models with optimized weights",
        "  • Grid search for best weight combination",
        "  • Model disagreement for confidence intervals"
    ], "Primary Models", "Ensemble Strategy")
    
    # Slide 7: Web Dashboard
    add_content_slide(prs, "Interactive Web Dashboard", [
        "Dashboard Page: Real-time sales trends, category breakdown, top products",
        "Product Catalog: Browse, search, filter by category, sort by metrics",
        "Predictions Page: Select product, choose forecast horizon (7-365 days)",
        "What-If Analysis: Simulate price changes, discounts, rating improvements",
        "Feature Importance: Understand which factors drive sales the most",
        "Training Interface: One-click model training with progress tracking",
        "Responsive design with hover states, transitions, micro-interactions",
        "Built with FastAPI + Jinja2 templates + Chart.js visualizations"
    ])
    
    # Slide 8: Model Training & Evaluation
    add_table_slide(prs, "Model Training & Evaluation", [
        "Metric", "Description", "Purpose"
    ], [
        ["RMSE", "Root Mean Square Error", "Overall prediction accuracy"],
        ["MAE", "Mean Absolute Error", "Average absolute prediction error"],
        ["MAPE", "Mean Absolute Percentage Error", "Relative error percentage"],
        ["R²", "Coefficient of Determination", "Variance explained by model"],
        ["Bias", "Systematic over/under prediction", "Directional accuracy"],
        ["Max Error", "Maximum absolute error", "Worst-case scenario"],
        ["Median Abs Error", "Median absolute error", "Typical error magnitude"]
    ])
    
    # Slide 9: Deployment Options
    add_two_column_slide(prs, "Deployment Options", [
        "Local Development:",
        "  • Python 3.9+ required",
        "  • pip install -r requirements.txt",
        "  • python run.py (launches dashboard)",
        "  • Automatic data generation",
        "  • Manual model training via UI",
        "",
        "Render (Production):",
        "  • Free tier deployment",
        "  • Lazy loading for fast startup",
        "  • Health check endpoint",
        "  • Pre-trained models included"
    ], [
        "Vercel (Serverless):",
        "  • Serverless entry point (index.py)",
        "  • Optimized for cold starts",
        "  • API-only (no dashboard)",
        "  • Automatic scaling",
        "",
        "Key Features:",
        "  • Background training threads",
        "  • Progress tracking API",
        "  • Model caching to disk",
        "  • Environment-aware configuration"
    ], "Running Locally", "Cloud Deployment")
    
    # Slide 10: Results & Demo - Screenshots
    if len(screenshot_files) >= 1:
        add_image_slide(prs, "Dashboard Overview", str(screenshot_files[0]), "Main dashboard showing sales trends and category breakdown")
    if len(screenshot_files) >= 2:
        add_image_slide(prs, "Product Analytics", str(screenshot_files[1]), "Detailed product analytics and performance metrics")
    if len(screenshot_files) >= 3:
        add_image_slide(prs, "Prediction Interface", str(screenshot_files[2]), "Sales prediction interface with confidence intervals")
    
    # If no screenshots, add a content slide
    if len(screenshot_files) == 0:
        add_content_slide(prs, "Results & Demo", [
            "Dashboard shows real-time sales metrics and visualizations",
            "Interactive charts for daily trends and category breakdown",
            "Product catalog with search, filter, and sort capabilities",
            "Prediction interface with confidence intervals",
            "What-if analysis for price/discount scenarios",
            "Feature importance viewer for model interpretability",
            "Training interface with progress tracking",
            "Responsive design works on desktop and mobile"
        ])
    
    # Slide 11: Future Improvements
    add_content_slide(prs, "Future Improvements", [
        "Real Data Integration: Connect to actual Amazon APIs or Kaggle datasets",
        "Advanced Models: Add LSTM/Transformer time-series models",
        "Real-time Updates: Live data streaming for up-to-date predictions",
        "User Authentication: Multi-user support with saved preferences",
        "API Documentation: Swagger/OpenAPI documentation for external integration",
        "Mobile App: React Native or Flutter mobile application",
        "A/B Testing: Framework for testing different prediction strategies",
        "Automated Retraining: Scheduled model retraining with new data"
    ])
    
    # Slide 12: Q&A
    add_title_slide(prs, "Questions & Discussion", "Thank you for your attention!")
    
    # Save presentation
    from datetime import datetime
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_path = Path(__file__).parent / f"Amazon_Sales_Prediction_System_0th_Review_{timestamp}.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    
    return output_path

if __name__ == "__main__":
    create_presentation()
